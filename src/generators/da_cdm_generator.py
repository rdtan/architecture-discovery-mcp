"""DA-CDM Conceptual Data Model diagram generator.

Produces a PPTX with entity nodes grouped by data domain,
connected by relationship lines labeled with cardinality.
"""

from pathlib import Path
import math
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from src.models.project import ProjectInfo, DataEntity, DataRelationship
from src.generators.pptx_generator import (
    create_presentation, add_rect_node, add_connector, save_presentation, COLORS,
)
from src.i18n import t


def generate_da_cdm(
    project: ProjectInfo,
    entities: list[DataEntity],
    relationships: list[DataRelationship],
    output_dir: Path,
    locale: str = "zh",
) -> Path:
    prs = create_presentation(t("pptx.cdm_diagram", locale, name=project.name))
    slide = prs.slides[0]

    if not entities:
        return save_presentation(prs, output_dir / t("file.da_cdm", locale))

    # Group entities by domain
    domain_order: list[str] = []
    domain_entities: dict[str, list[DataEntity]] = {}
    for entity in entities:
        domain = entity.data_domain or entity.module_name
        if domain not in domain_entities:
            domain_order.append(domain)
            domain_entities[domain] = []
        domain_entities[domain].append(entity)

    # Layout parameters
    page_left = Inches(0.4)
    page_top = Inches(1.1)
    page_width = Inches(12.5)
    page_height = Inches(5.8)

    node_width = Inches(1.6)
    node_height = Inches(0.7)
    node_gap_x = Inches(0.4)
    node_gap_y = Inches(0.5)
    domain_gap_y = Inches(0.8)
    domain_title_height = Inches(0.35)
    max_cols = 5

    # Place entity nodes, track positions for connectors
    node_positions: dict[str, tuple] = {}  # class_name -> (center_x, center_y)
    current_y = page_top

    for domain_idx, domain in enumerate(domain_order):
        ents = domain_entities[domain]
        color = COLORS[domain_idx % len(COLORS)]

        # Domain title
        title_box = slide.shapes.add_textbox(
            page_left, current_y, page_width, domain_title_height
        )
        tf = title_box.text_frame
        tf.paragraphs[0].text = t("pptx.domain_group", locale, name=domain)
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = color

        current_y += domain_title_height + Inches(0.1)

        # Place entities in grid
        for idx, entity in enumerate(ents):
            col = idx % max_cols
            row = idx // max_cols

            x = page_left + (node_width + node_gap_x) * col
            y = current_y + (node_height + node_gap_y) * row

            # Skip if off page
            if y + node_height > page_top + page_height:
                break

            label = entity.class_name
            add_rect_node(slide, x, y, node_width, node_height, label, color, font_size=Pt(8))

            node_positions[entity.class_name] = (
                x + node_width / 2,
                y + node_height / 2,
            )

        # Advance Y for next domain
        num_rows = (len(ents) + max_cols - 1) // max_cols
        current_y += (node_height + node_gap_y) * num_rows + domain_gap_y

        if current_y > page_top + page_height:
            break

    # Draw relationship lines
    for rel in relationships:
        if rel.source_entity in node_positions and rel.target_entity in node_positions:
            sx, sy = node_positions[rel.source_entity]
            tx, ty = node_positions[rel.target_entity]
            add_connector(slide, sx, sy, tx, ty)

            # Label at midpoint
            label_x = (sx + tx) / 2 - Inches(0.4)
            label_y = (sy + ty) / 2 - Inches(0.15)
            txBox = slide.shapes.add_textbox(label_x, label_y, Inches(0.8), Inches(0.25))
            tf = txBox.text_frame
            tf.paragraphs[0].text = rel.relationship_type
            tf.paragraphs[0].font.size = Pt(7)
            tf.paragraphs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Legend
    _draw_relationship_legend(slide, locale)

    return save_presentation(prs, output_dir / t("file.da_cdm", locale))


def _draw_relationship_legend(slide, locale: str) -> None:
    legend_top = Inches(6.8)
    legend_left = Inches(0.5)

    title_box = slide.shapes.add_textbox(legend_left, legend_top, Inches(1.5), Inches(0.25))
    tf = title_box.text_frame
    tf.paragraphs[0].text = t("pptx.legend_relationship", locale)
    tf.paragraphs[0].font.size = Pt(8)
    tf.paragraphs[0].font.bold = True

    rel_types = ["1:1", "1:N", "N:1", "N:N"]
    for idx, rel_type in enumerate(rel_types):
        x = legend_left + Inches(1.8) + Inches(idx * 1.2)
        txBox = slide.shapes.add_textbox(x, legend_top, Inches(1.0), Inches(0.25))
        tf = txBox.text_frame
        tf.paragraphs[0].text = f"— {rel_type}"
        tf.paragraphs[0].font.size = Pt(7)
        tf.paragraphs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)
