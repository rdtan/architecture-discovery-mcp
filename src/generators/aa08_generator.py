from pathlib import Path
import math
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from src.models.project import ProjectInfo, Integration
from src.generators.pptx_generator import (
    create_presentation, add_rect_node, add_connector, save_presentation, COLORS,
)
from src.i18n import t

INTEGRATION_COLORS = {
    "HTTP": RGBColor(0x44, 0x72, 0xC4),
    "RPC": RGBColor(0xED, 0x7D, 0x31),
    "MQ": RGBColor(0x70, 0xAD, 0x47),
}


def generate_aa08(project: ProjectInfo, integrations: list[Integration], output_dir: Path, locale: str = "zh") -> Path:
    prs = create_presentation(t("pptx.integration_diagram", locale, name=project.name))
    slide = prs.slides[0]

    if not integrations:
        return save_presentation(prs, output_dir / t("file.aa08", locale))

    nodes = set()
    for intg in integrations:
        nodes.add(intg.source_module)
        nodes.add(intg.target_module)
    nodes = sorted(nodes)

    center_x = Inches(6.5)
    center_y = Inches(4.0)
    radius = Inches(2.5)
    node_width = Inches(1.8)
    node_height = Inches(0.8)

    node_positions = {}
    for idx, node_name in enumerate(nodes):
        angle = (2 * math.pi * idx) / len(nodes) - math.pi / 2
        x = center_x + radius * math.cos(angle) - node_width / 2
        y = center_y + radius * math.sin(angle) - node_height / 2

        color = COLORS[idx % len(COLORS)]
        add_rect_node(slide, x, y, node_width, node_height, node_name, color, font_size=Pt(8))
        node_positions[node_name] = (
            x + node_width / 2,
            y + node_height / 2,
        )

    for intg in integrations:
        if intg.source_module in node_positions and intg.target_module in node_positions:
            sx, sy = node_positions[intg.source_module]
            tx, ty = node_positions[intg.target_module]
            add_connector(slide, sx, sy, tx, ty)

            label_x = (sx + tx) / 2 - Inches(0.5)
            label_y = (sy + ty) / 2 - Inches(0.2)
            txBox = slide.shapes.add_textbox(label_x, label_y, Inches(1.2), Inches(0.3))
            tf = txBox.text_frame
            tf.paragraphs[0].text = intg.integration_type.value
            tf.paragraphs[0].font.size = Pt(7)
            color = INTEGRATION_COLORS.get(intg.integration_type.value, RGBColor(0x66, 0x66, 0x66))
            tf.paragraphs[0].font.color.rgb = color

    legend_top = Inches(6.5)
    legend_left = Inches(0.5)
    for idx, (label, color) in enumerate(INTEGRATION_COLORS.items()):
        txBox = slide.shapes.add_textbox(
            legend_left + Inches(idx * 2), legend_top,
            Inches(1.5), Inches(0.3),
        )
        tf = txBox.text_frame
        tf.paragraphs[0].text = f"■ {label}"
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].font.color.rgb = color

    return save_presentation(prs, output_dir / t("file.aa08", locale))
