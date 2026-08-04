"""DA-LINEAGE Data Lineage Diagram generator.

Produces a PPTX with entity-level lineage flow diagrams.
Smart compact layout: merges small systems, caps large systems to top-N by degree,
uses multi-column grid layout for density.
"""

from collections import defaultdict
from pathlib import Path

import networkx as nx
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

from src.models.project import ProjectInfo
from src.analyzers.lineage_graph import LineageGraph
from src.generators.pptx_generator import (
    create_presentation, add_rect_node, add_connector, save_presentation, COLORS,
)
from src.i18n import t

# Thresholds for smart partitioning
LARGE_SYSTEM_THRESHOLD = 30
TOP_N_ENTITIES = 25
SMALL_SYSTEM_THRESHOLD = 10

# Compact layout constants
MARGIN_LEFT = Inches(0.3)
MARGIN_TOP = Inches(0.85)
USABLE_WIDTH = Inches(12.7)
USABLE_HEIGHT = Inches(6.2)

ENTITY_BOX_WIDTH = Inches(1.3)
ENTITY_BOX_HEIGHT = Inches(0.35)
ENTITY_GAP_X = Inches(0.1)
ENTITY_GAP_Y = Inches(0.08)
SYSTEM_HEADER_HEIGHT = Inches(0.24)
SYSTEM_GAP_Y = Inches(0.15)
LANE_PADDING = Inches(0.08)

MAPPING_TYPE_COLORS = {
    "direct": RGBColor(0x44, 0x72, 0xC4),
    "transform": RGBColor(0xED, 0x7D, 0x31),
    "aggregate": RGBColor(0x70, 0xAD, 0x47),
    "concat": RGBColor(0xFF, 0xC0, 0x00),
    "constant": RGBColor(0xA5, 0xA5, 0xA5),
    "flow": RGBColor(0x5B, 0x9B, 0xD5),
}


def _stable_color_index(name: str) -> int:
    """Deterministic color index based on string content (not Python hash())."""
    h = 0
    for ch in name:
        h = (h * 31 + ord(ch)) & 0xFFFFFFFF
    return h % len(COLORS)


def generate_da_lineage(
    project: ProjectInfo,
    graph: LineageGraph,
    output_dir: Path,
    locale: str = "zh",
) -> Path:
    """Generate DA-LINEAGE data lineage PPTX diagram."""
    title = t("pptx.lineage_diagram", locale, name=project.name)
    prs = create_presentation(title)

    if graph.node_count == 0:
        return save_presentation(prs, output_dir / t("file.da_lineage", locale))

    systems, entity_edges, field_counts = _build_entity_graph(graph)

    # Compute entity degrees for prioritization
    entity_degree = _compute_entity_degrees(entity_edges)

    # Smart partition into slides
    slide_groups = _smart_partition(systems, entity_edges, entity_degree)

    # Render first slide (reuse the title slide from create_presentation)
    first_group = slide_groups[0]
    _render_compact_slide(
        prs.slides[0], first_group, systems, entity_edges,
        field_counts, entity_degree, locale
    )

    # Additional slides if needed
    for group in slide_groups[1:]:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        _render_compact_slide(
            slide, group, systems, entity_edges,
            field_counts, entity_degree, locale
        )

    return save_presentation(prs, output_dir / t("file.da_lineage", locale))


def _build_entity_graph(graph: LineageGraph):
    """Collapse field-level graph to entity-level."""
    nodes = graph.get_all_nodes()
    edges = graph.get_all_edges()

    systems_set: dict[str, set] = defaultdict(set)
    field_counts: dict[str, set] = defaultdict(set)

    for node_key in nodes:
        parts = node_key.split(".", 2)
        if len(parts) >= 3:
            sys_name, entity_name, field_name = parts[0], parts[1], parts[2]
            systems_set[sys_name].add(entity_name)
            entity_key = f"{sys_name}.{entity_name}"
            field_counts[entity_key].add(field_name)

    systems = {k: sorted(v) for k, v in systems_set.items()}
    field_count_ints = {k: len(v) for k, v in field_counts.items()}

    entity_edges: dict[tuple, dict] = defaultdict(
        lambda: {"count": 0, "mapping_types": set(), "type_counts": defaultdict(int)}
    )

    for u, v, data in edges:
        u_parts = u.split(".", 2)
        v_parts = v.split(".", 2)
        if len(u_parts) < 2 or len(v_parts) < 2:
            continue

        src_entity = f"{u_parts[0]}.{u_parts[1]}"
        tgt_entity = f"{v_parts[0]}.{v_parts[1]}"

        if src_entity == tgt_entity:
            continue

        key = (src_entity, tgt_entity)
        mt = data.get("mapping_type", "direct")
        entity_edges[key]["count"] += 1
        entity_edges[key]["mapping_types"].add(mt)
        entity_edges[key]["type_counts"][mt] += 1

    return systems, dict(entity_edges), field_count_ints


def _compute_entity_degrees(entity_edges: dict) -> dict[str, int]:
    """Compute degree (in + out) for each entity."""
    degree: dict[str, int] = defaultdict(int)
    for (src, tgt) in entity_edges:
        degree[src] += 1
        degree[tgt] += 1
    return dict(degree)


def _smart_partition(systems: dict, entity_edges: dict, entity_degree: dict) -> list[dict]:
    """Smart partition: large systems get own slide with top-N, small systems merge."""
    # Categorize systems
    large_systems = []
    medium_systems = []
    small_systems = []

    for sys_name, entities in sorted(systems.items(), key=lambda x: -len(x[1])):
        if len(entities) >= LARGE_SYSTEM_THRESHOLD:
            large_systems.append(sys_name)
        elif len(entities) >= SMALL_SYSTEM_THRESHOLD:
            medium_systems.append(sys_name)
        else:
            small_systems.append(sys_name)

    slide_groups = []

    # Each large system gets its own slide with top-N entities by degree
    for sys_name in large_systems:
        entities = systems[sys_name]
        entity_keys = [f"{sys_name}.{e}" for e in entities]
        # Keep only entities that participate in edges, sorted by degree
        active_keys = [ek for ek in entity_keys if entity_degree.get(ek, 0) > 0]
        active_keys.sort(key=lambda ek: entity_degree.get(ek, 0), reverse=True)
        top_keys = set(active_keys[:TOP_N_ENTITIES])
        overflow = len(active_keys) - len(top_keys)

        slide_groups.append({
            "title": sys_name,
            "systems_included": [sys_name],
            "entities": top_keys,
            "overflow_count": overflow,
        })

    # Medium systems merged into one slide
    if medium_systems:
        merged_entities = set()
        for sys_name in medium_systems:
            entity_keys = [f"{sys_name}.{e}" for e in systems[sys_name]]
            active_keys = [ek for ek in entity_keys if entity_degree.get(ek, 0) > 0]
            active_keys.sort(key=lambda ek: entity_degree.get(ek, 0), reverse=True)
            merged_entities.update(active_keys[:TOP_N_ENTITIES])

        slide_groups.append({
            "title": " + ".join(medium_systems),
            "systems_included": medium_systems,
            "entities": merged_entities,
            "overflow_count": 0,
        })

    # Small systems all merged into one slide
    if small_systems:
        merged_entities = set()
        for sys_name in small_systems:
            entity_keys = [f"{sys_name}.{e}" for e in systems[sys_name]]
            active_keys = [ek for ek in entity_keys if entity_degree.get(ek, 0) > 0]
            merged_entities.update(active_keys)

        if merged_entities:
            slide_groups.append({
                "title": "Others",
                "systems_included": small_systems,
                "entities": merged_entities,
                "overflow_count": 0,
            })

    return slide_groups if slide_groups else [{"title": "", "systems_included": list(systems.keys()), "entities": set(), "overflow_count": 0}]


def _render_compact_slide(slide, group, systems, entity_edges, field_counts, entity_degree, locale):
    """Render a compact slide with multi-column grid layout."""
    title_text = t("pptx.lineage_diagram", locale, name=group["title"])
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12), Inches(0.45))
    tf = txBox.text_frame
    tf.paragraphs[0].text = title_text
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    included_entities = group["entities"]
    if not included_entities:
        return

    # Group entities by system for layout
    sys_entities: dict[str, list[str]] = defaultdict(list)
    for ek in included_entities:
        sys_name = ek.split(".")[0]
        sys_entities[sys_name].append(ek)

    # Sort entities within each system by degree (most connected first)
    for sys_name in sys_entities:
        sys_entities[sys_name].sort(key=lambda ek: entity_degree.get(ek, 0), reverse=True)

    # Multi-column grid layout
    positions, overflow_indicators = _grid_layout(sys_entities, group.get("overflow_count", 0))

    # Draw system lane backgrounds
    _draw_compact_lanes(slide, sys_entities, positions)

    # Draw entity boxes
    centers: dict[str, tuple] = {}
    for entity_key, (left, top) in positions.items():
        sys_name = entity_key.split(".")[0]
        entity_name = entity_key.split(".", 1)[1]
        fc = field_counts.get(entity_key, 0)
        color = COLORS[_stable_color_index(sys_name)]

        label = f"{entity_name}"
        if fc > 0:
            label += f" ({fc})"

        add_rect_node(slide, left, top, ENTITY_BOX_WIDTH, ENTITY_BOX_HEIGHT, label, color, Pt(6))
        centers[entity_key] = (left + ENTITY_BOX_WIDTH / 2, top + ENTITY_BOX_HEIGHT / 2)

    # Draw overflow indicators
    for (ox, oy, sys_name, count) in overflow_indicators:
        label = t("pptx.lineage_more_entities", locale, count=count)
        txBox = slide.shapes.add_textbox(int(ox), int(oy), int(ENTITY_BOX_WIDTH), int(Inches(0.2)))
        tf = txBox.text_frame
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.size = Pt(5)
        tf.paragraphs[0].font.italic = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    # Draw edges
    slide_mapping_types = set()
    edge_offset_tracker: dict[str, int] = defaultdict(int)
    for (src_ent, tgt_ent), info in entity_edges.items():
        if src_ent in centers and tgt_ent in centers:
            sx, sy = centers[src_ent]
            tx, ty = centers[tgt_ent]

            dominant_type = max(info["mapping_types"], key=lambda mt: info.get("type_counts", {}).get(mt, 0)) if info["mapping_types"] else "direct"
            color = MAPPING_TYPE_COLORS.get(dominant_type, RGBColor(0x66, 0x66, 0x66))
            slide_mapping_types.add(dominant_type)

            src_out_count = edge_offset_tracker[f"out:{src_ent}"]
            edge_offset_tracker[f"out:{src_ent}"] += 1
            y_offset = Inches(0.03) * (src_out_count % 3 - 1)

            connector = slide.shapes.add_connector(
                1, int(sx), int(sy + y_offset), int(tx), int(ty + y_offset)
            )
            connector.line.color.rgb = color
            connector.line.width = Pt(0.75)
            _add_arrowhead(connector)

    # Legend
    _draw_legend(slide, locale, slide_mapping_types)


def _grid_layout(sys_entities: dict[str, list[str]], total_overflow: int):
    """Assign positions using a multi-column grid per system.

    Systems are laid out left-to-right. Within each system column,
    entities fill top-to-bottom, wrapping to a new column when vertical space runs out.
    """
    positions: dict[str, tuple] = {}
    overflow_indicators: list[tuple] = []

    # Determine how many system "columns" we need
    num_systems = len(sys_entities)
    if num_systems == 0:
        return positions, overflow_indicators

    # Calculate how many entity columns each system gets based on entity count
    total_entities = sum(len(ents) for ents in sys_entities.values())
    max_rows_per_col = int((USABLE_HEIGHT - SYSTEM_HEADER_HEIGHT) / (ENTITY_BOX_HEIGHT + ENTITY_GAP_Y))

    # Allocate horizontal space proportionally
    system_allocs = []
    for sys_name, ents in sorted(sys_entities.items(), key=lambda x: -len(x[1])):
        n_ents = len(ents)
        cols_needed = max(1, (n_ents + max_rows_per_col - 1) // max_rows_per_col)
        system_allocs.append((sys_name, ents, cols_needed))

    total_cols = sum(a[2] for a in system_allocs)
    col_width = ENTITY_BOX_WIDTH + ENTITY_GAP_X

    # If total columns exceed available width, cap per system
    max_total_cols = int(USABLE_WIDTH / col_width)
    if total_cols > max_total_cols:
        # Redistribute: cap each system proportionally
        scale = max_total_cols / total_cols
        for i, (sys_name, ents, cols) in enumerate(system_allocs):
            system_allocs[i] = (sys_name, ents, max(1, int(cols * scale)))
        total_cols = sum(a[2] for a in system_allocs)

    # Position each system's entities
    current_x = MARGIN_LEFT

    for sys_name, ents, cols_for_sys in system_allocs:
        sys_width = cols_for_sys * col_width
        base_y = MARGIN_TOP + SYSTEM_HEADER_HEIGHT + Inches(0.05)

        placed = 0
        for col_idx in range(cols_for_sys):
            col_x = current_x + col_idx * col_width
            row_y = base_y

            while placed < len(ents):
                if row_y + ENTITY_BOX_HEIGHT > MARGIN_TOP + USABLE_HEIGHT:
                    break
                entity_key = ents[placed]
                positions[entity_key] = (col_x, row_y)
                row_y += ENTITY_BOX_HEIGHT + ENTITY_GAP_Y
                placed += 1

        if placed < len(ents):
            overflow_count = len(ents) - placed + total_overflow
            overflow_indicators.append((current_x, MARGIN_TOP + USABLE_HEIGHT - Inches(0.2), sys_name, overflow_count))
        elif total_overflow > 0 and sys_name == system_allocs[0][0]:
            overflow_indicators.append((current_x, MARGIN_TOP + USABLE_HEIGHT - Inches(0.2), sys_name, total_overflow))

        current_x += sys_width + Inches(0.15)

    return positions, overflow_indicators


def _draw_compact_lanes(slide, sys_entities, positions):
    """Draw compact system lane backgrounds."""
    if not positions:
        return

    # Group positions by system
    sys_bounds: dict[str, dict] = {}
    for entity_key, (left, top) in positions.items():
        sys_name = entity_key.split(".")[0]
        if sys_name not in sys_bounds:
            sys_bounds[sys_name] = {"min_x": left, "max_x": left, "min_y": top, "max_y": top}
        else:
            b = sys_bounds[sys_name]
            b["min_x"] = min(b["min_x"], left)
            b["max_x"] = max(b["max_x"], left)
            b["min_y"] = min(b["min_y"], top)
            b["max_y"] = max(b["max_y"], top)

    for sys_name, bounds in sys_bounds.items():
        lane_left = bounds["min_x"] - LANE_PADDING
        lane_top = bounds["min_y"] - SYSTEM_HEADER_HEIGHT - Inches(0.05)
        lane_width = bounds["max_x"] - bounds["min_x"] + ENTITY_BOX_WIDTH + LANE_PADDING * 2
        lane_height = bounds["max_y"] - bounds["min_y"] + ENTITY_BOX_HEIGHT + SYSTEM_HEADER_HEIGHT + Inches(0.15)

        # Background
        bg = slide.shapes.add_shape(
            1, int(lane_left), int(lane_top),
            int(lane_width), int(lane_height)
        )
        base_color = COLORS[_stable_color_index(sys_name)]
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(
            min(255, base_color[0] + 200),
            min(255, base_color[1] + 200),
            min(255, base_color[2] + 200),
        )
        bg.line.fill.background()

        # System label
        txBox = slide.shapes.add_textbox(
            int(lane_left + Inches(0.03)), int(lane_top + Inches(0.02)),
            int(lane_width - Inches(0.06)), int(SYSTEM_HEADER_HEIGHT)
        )
        tf = txBox.text_frame
        tf.paragraphs[0].text = sys_name
        tf.paragraphs[0].font.size = Pt(6)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = base_color


def _draw_legend(slide, locale, visible_types=None):
    """Draw mapping type legend at the bottom, filtered to visible types."""
    legend_y = Inches(7.1)
    legend_x = Inches(0.3)

    title_text = t("pptx.lineage_legend_title", locale)
    txBox = slide.shapes.add_textbox(int(legend_x), int(legend_y), Inches(0.8), Inches(0.18))
    tf = txBox.text_frame
    tf.paragraphs[0].text = title_text
    tf.paragraphs[0].font.size = Pt(6)
    tf.paragraphs[0].font.bold = True

    x_offset = legend_x + Inches(0.9)
    for mt, color in MAPPING_TYPE_COLORS.items():
        if visible_types is not None and mt not in visible_types:
            continue
        box = slide.shapes.add_shape(1, int(x_offset), int(legend_y + Inches(0.02)), Inches(0.12), Inches(0.12))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        txBox = slide.shapes.add_textbox(int(x_offset + Inches(0.15)), int(legend_y), Inches(0.7), Inches(0.18))
        tf = txBox.text_frame
        tf.paragraphs[0].text = mt
        tf.paragraphs[0].font.size = Pt(5)

        x_offset += Inches(1.1)


def _add_arrowhead(connector):
    """Add a triangle arrowhead to the end of a connector via OOXML."""
    ln = connector.line._ln
    tail_end = ln.makeelement(qn('a:tailEnd'), {})
    tail_end.set('type', 'triangle')
    tail_end.set('w', 'med')
    tail_end.set('len', 'med')
    ln.append(tail_end)
