"""DA-LDM Logical Data Model diagram generator.

Produces a PPTX with detailed entity cards showing fields,
grouped by data domain (one domain per slide if many entities).
Relationship lines connect entities with cardinality and FK labels.
"""

from pathlib import Path
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from src.models.project import ProjectInfo, DataEntity, DataRelationship
from src.generators.pptx_generator import (
    create_presentation, add_connector, save_presentation, COLORS,
)
from src.i18n import t


MAX_FIELDS_DISPLAY = 8
NODE_WIDTH = Inches(2.0)
TITLE_HEIGHT = Inches(0.3)
FIELD_ROW_HEIGHT = Inches(0.22)
NODE_GAP_X = Inches(0.5)
NODE_GAP_Y = Inches(0.4)
MAX_COLS = 5


def generate_da_ldm(
    project: ProjectInfo,
    entities: list[DataEntity],
    relationships: list[DataRelationship],
    output_dir: Path,
    locale: str = "zh",
) -> Path:
    prs = create_presentation(t("pptx.ldm_diagram", locale, name=project.name))

    if not entities:
        return save_presentation(prs, output_dir / t("file.da_ldm", locale))

    # Group entities by domain
    domain_order: list[str] = []
    domain_entities: dict[str, list[DataEntity]] = {}
    for entity in entities:
        domain = entity.data_domain or entity.module_name
        if domain not in domain_entities:
            domain_order.append(domain)
            domain_entities[domain] = []
        domain_entities[domain].append(entity)

    # Decide: if total entities <= 12, single slide; otherwise per-domain slides
    total_entities = len(entities)
    use_single_slide = total_entities <= 12

    if use_single_slide:
        slide = prs.slides[0]
        node_positions = _draw_entities_on_slide(slide, entities, domain_order, domain_entities, locale)
        _draw_relationships(slide, relationships, node_positions)
    else:
        # Remove default empty slide, create per-domain slides
        slide = prs.slides[0]
        all_node_positions: dict[str, tuple] = {}

        # Use first slide for first domain
        for domain_idx, domain in enumerate(domain_order):
            ents = domain_entities[domain]
            if domain_idx == 0:
                current_slide = slide
            else:
                slide_layout = prs.slide_layouts[6]
                current_slide = prs.slides.add_slide(slide_layout)
                # Add title
                txBox = current_slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.3), Inches(12), Inches(0.6)
                )
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = t("pptx.ldm_diagram", locale, name=project.name)
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

            positions = _draw_domain_slide(current_slide, domain, ents, domain_idx, locale)
            all_node_positions.update(positions)

            # Draw relationships for entities on this slide
            domain_class_names = {e.class_name for e in ents}
            domain_rels = [
                r for r in relationships
                if r.source_entity in domain_class_names or r.target_entity in domain_class_names
            ]
            _draw_relationships(current_slide, domain_rels, positions)

    return save_presentation(prs, output_dir / t("file.da_ldm", locale))


def _draw_entities_on_slide(
    slide, entities: list[DataEntity],
    domain_order: list[str],
    domain_entities: dict[str, list[DataEntity]],
    locale: str,
) -> dict[str, tuple]:
    """Draw all entities on a single slide, grouped by domain."""
    node_positions: dict[str, tuple] = {}
    page_left = Inches(0.4)
    page_top = Inches(1.1)

    current_y = page_top

    for domain_idx, domain in enumerate(domain_order):
        ents = domain_entities[domain]
        color = COLORS[domain_idx % len(COLORS)]

        # Domain label
        title_box = slide.shapes.add_textbox(page_left, current_y, Inches(12), Inches(0.3))
        tf = title_box.text_frame
        tf.paragraphs[0].text = t("pptx.domain_group", locale, name=domain)
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = color

        current_y += Inches(0.35)

        for idx, entity in enumerate(ents):
            col = idx % MAX_COLS
            row = idx // MAX_COLS

            node_height = _calc_node_height(entity)
            x = page_left + (NODE_WIDTH + NODE_GAP_X) * col
            y = current_y + (node_height + NODE_GAP_Y) * row

            if y + node_height > Inches(7.2):
                break

            _draw_entity_card(slide, x, y, entity, color)
            node_positions[entity.class_name] = (
                x + NODE_WIDTH / 2,
                y + node_height / 2,
            )

        num_rows = (len(ents) + MAX_COLS - 1) // MAX_COLS
        tallest = max((_calc_node_height(e) for e in ents), default=Inches(1.0))
        current_y += (tallest + NODE_GAP_Y) * num_rows + Inches(0.3)

    return node_positions


def _draw_domain_slide(
    slide, domain: str, entities: list[DataEntity], domain_idx: int, locale: str
) -> dict[str, tuple]:
    """Draw entities for a single domain on its dedicated slide."""
    node_positions: dict[str, tuple] = {}
    color = COLORS[domain_idx % len(COLORS)]
    page_left = Inches(0.4)
    page_top = Inches(1.1)

    # Domain label
    title_box = slide.shapes.add_textbox(page_left, page_top, Inches(12), Inches(0.3))
    tf = title_box.text_frame
    tf.paragraphs[0].text = t("pptx.domain_group", locale, name=domain)
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = color

    start_y = page_top + Inches(0.4)

    for idx, entity in enumerate(entities):
        col = idx % MAX_COLS
        row = idx // MAX_COLS

        node_height = _calc_node_height(entity)
        x = page_left + (NODE_WIDTH + NODE_GAP_X) * col
        y = start_y + (node_height + NODE_GAP_Y) * row

        if y + node_height > Inches(7.2):
            break

        _draw_entity_card(slide, x, y, entity, color)
        node_positions[entity.class_name] = (
            x + NODE_WIDTH / 2,
            y + node_height / 2,
        )

    return node_positions


def _calc_node_height(entity: DataEntity) -> int:
    """Calculate node height based on field count."""
    field_count = min(len(entity.fields), MAX_FIELDS_DISPLAY)
    overflow = 1 if len(entity.fields) > MAX_FIELDS_DISPLAY else 0
    return TITLE_HEIGHT + FIELD_ROW_HEIGHT * (field_count + overflow) + Inches(0.1)


def _draw_entity_card(slide, left, top, entity: DataEntity, color: RGBColor) -> None:
    """Draw an entity card with title bar and field list."""
    node_height = _calc_node_height(entity)

    # Background container
    bg_shape = slide.shapes.add_shape(1, left, top, NODE_WIDTH, node_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    bg_shape.line.color.rgb = color
    bg_shape.line.width = Pt(1.0)

    # Title bar
    title_shape = slide.shapes.add_shape(1, left, top, NODE_WIDTH, TITLE_HEIGHT)
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = color
    title_shape.line.color.rgb = color
    tf = title_shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = entity.class_name
    tf.paragraphs[0].font.size = Pt(7)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Field rows
    fields_to_show = entity.fields[:MAX_FIELDS_DISPLAY]
    for idx, field in enumerate(fields_to_show):
        y = top + TITLE_HEIGHT + FIELD_ROW_HEIGHT * idx + Inches(0.04)

        prefix = ""
        if field.is_primary_key:
            prefix = "PK "
        elif field.is_foreign_key:
            prefix = "FK "

        text = f"{prefix}{field.name}: {field.java_type}"
        if len(text) > 24:
            text = text[:22] + ".."

        txBox = slide.shapes.add_textbox(left + Inches(0.05), y, NODE_WIDTH - Inches(0.1), FIELD_ROW_HEIGHT)
        tf = txBox.text_frame
        tf.word_wrap = False
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(6)
        tf.paragraphs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        if field.is_primary_key:
            tf.paragraphs[0].font.bold = True

    # Overflow indicator
    if len(entity.fields) > MAX_FIELDS_DISPLAY:
        y = top + TITLE_HEIGHT + FIELD_ROW_HEIGHT * MAX_FIELDS_DISPLAY + Inches(0.04)
        remaining = len(entity.fields) - MAX_FIELDS_DISPLAY
        txBox = slide.shapes.add_textbox(left + Inches(0.05), y, NODE_WIDTH - Inches(0.1), FIELD_ROW_HEIGHT)
        tf = txBox.text_frame
        tf.paragraphs[0].text = f"... +{remaining} fields"
        tf.paragraphs[0].font.size = Pt(6)
        tf.paragraphs[0].font.color.rgb = RGBColor(0x99, 0x99, 0x99)


def _draw_relationships(
    slide, relationships: list[DataRelationship], node_positions: dict[str, tuple]
) -> None:
    """Draw relationship connector lines between entities."""
    for rel in relationships:
        if rel.source_entity in node_positions and rel.target_entity in node_positions:
            sx, sy = node_positions[rel.source_entity]
            tx, ty = node_positions[rel.target_entity]
            add_connector(slide, sx, sy, tx, ty)

            # Label at midpoint
            label_x = (sx + tx) / 2 - Inches(0.5)
            label_y = (sy + ty) / 2 - Inches(0.15)
            label_text = rel.relationship_type
            if rel.fk_field:
                label_text += f" ({rel.fk_field})"

            txBox = slide.shapes.add_textbox(label_x, label_y, Inches(1.2), Inches(0.25))
            tf = txBox.text_frame
            tf.paragraphs[0].text = label_text
            tf.paragraphs[0].font.size = Pt(6)
            tf.paragraphs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
