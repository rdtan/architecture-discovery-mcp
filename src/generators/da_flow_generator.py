"""DA-FLOW Data Flow diagram generator.

Produces a PPTX showing which application functions perform CRUD
operations on which data entities. Entities are placed in the center,
application modules around them, with colored connectors indicating
operation type (C/R/U/D).
"""

from pathlib import Path
import math
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from src.models.project import ProjectInfo, DataEntity
from src.generators.pptx_generator import (
    create_presentation, add_rect_node, add_connector, save_presentation, COLORS,
)
from src.i18n import t


CRUD_COLORS = {
    "C": RGBColor(0x70, 0xAD, 0x47),  # Green
    "R": RGBColor(0x44, 0x72, 0xC4),  # Blue
    "U": RGBColor(0xED, 0x7D, 0x31),  # Orange
    "D": RGBColor(0xC0, 0x39, 0x2B),  # Red
}

CRUD_LABELS = {
    "C": "Create",
    "R": "Read",
    "U": "Update",
    "D": "Delete",
}


def generate_da_flow(
    project: ProjectInfo,
    entities: list[DataEntity],
    crud_records: list[dict],
    output_dir: Path,
    locale: str = "zh",
) -> Path:
    prs = create_presentation(t("pptx.data_flow_diagram", locale, name=project.name))
    slide = prs.slides[0]

    if not entities or not crud_records:
        return save_presentation(prs, output_dir / t("file.da_flow", locale))

    # Collect unique app modules and entities involved in CRUD
    app_modules: set[str] = set()
    entity_names: set[str] = set()
    for record in crud_records:
        app_modules.add(record.get("app_name", ""))
        entity_names.add(record.get("entity", ""))

    app_modules.discard("")
    entity_names.discard("")

    if not app_modules or not entity_names:
        return save_presentation(prs, output_dir / t("file.da_flow", locale))

    # Layout: entities in center grid, app modules around the perimeter
    entity_positions = _place_entities_center(slide, sorted(entity_names))
    app_positions = _place_apps_perimeter(slide, sorted(app_modules))

    # Draw CRUD connections
    for record in crud_records:
        app_name = record.get("app_name", "")
        entity_name = record.get("entity", "")
        operation = record.get("operation", "")

        if app_name in app_positions and entity_name in entity_positions:
            sx, sy = app_positions[app_name]
            tx, ty = entity_positions[entity_name]

            color = CRUD_COLORS.get(operation, RGBColor(0x99, 0x99, 0x99))
            connector = slide.shapes.add_connector(
                1, int(sx), int(sy), int(tx), int(ty)
            )
            connector.line.color.rgb = color
            connector.line.width = Pt(1.0)

    # Legend
    _draw_crud_legend(slide, locale)

    return save_presentation(prs, output_dir / t("file.da_flow", locale))


def _place_entities_center(slide, entity_names: list[str]) -> dict[str, tuple]:
    """Place entity nodes in a center grid."""
    positions: dict[str, tuple] = {}

    center_x = Inches(6.5)
    center_y = Inches(3.8)
    node_width = Inches(1.5)
    node_height = Inches(0.6)
    gap_x = Inches(0.3)
    gap_y = Inches(0.3)

    max_cols = min(4, len(entity_names))
    if max_cols == 0:
        return positions

    num_rows = (len(entity_names) + max_cols - 1) // max_cols
    grid_width = max_cols * node_width + (max_cols - 1) * gap_x
    grid_height = num_rows * node_height + (num_rows - 1) * gap_y
    start_x = center_x - grid_width / 2
    start_y = center_y - grid_height / 2

    for idx, name in enumerate(entity_names):
        col = idx % max_cols
        row = idx // max_cols

        x = start_x + (node_width + gap_x) * col
        y = start_y + (node_height + gap_y) * row

        color = COLORS[1]  # Orange for entities
        add_rect_node(slide, x, y, node_width, node_height, name, color, font_size=Pt(7))

        positions[name] = (x + node_width / 2, y + node_height / 2)

    return positions


def _place_apps_perimeter(slide, app_names: list[str]) -> dict[str, tuple]:
    """Place application module nodes in a circle around the center."""
    positions: dict[str, tuple] = {}

    center_x = Inches(6.5)
    center_y = Inches(3.8)
    radius = Inches(3.0)
    node_width = Inches(1.6)
    node_height = Inches(0.6)

    for idx, name in enumerate(app_names):
        angle = (2 * math.pi * idx) / len(app_names) - math.pi / 2
        x = center_x + radius * math.cos(angle) - node_width / 2
        y = center_y + radius * math.sin(angle) - node_height / 2

        color = COLORS[0]  # Blue for apps
        add_rect_node(slide, x, y, node_width, node_height, name, color, font_size=Pt(7))

        positions[name] = (x + node_width / 2, y + node_height / 2)

    return positions


def _draw_crud_legend(slide, locale: str) -> None:
    """Draw color-coded CRUD operation legend at bottom."""
    legend_top = Inches(6.8)
    legend_left = Inches(0.5)

    title_box = slide.shapes.add_textbox(legend_left, legend_top, Inches(1.5), Inches(0.25))
    tf = title_box.text_frame
    tf.paragraphs[0].text = t("pptx.legend_crud", locale)
    tf.paragraphs[0].font.size = Pt(8)
    tf.paragraphs[0].font.bold = True

    for idx, (op, color) in enumerate(CRUD_COLORS.items()):
        label = CRUD_LABELS[op]
        x = legend_left + Inches(1.8) + Inches(idx * 1.5)
        txBox = slide.shapes.add_textbox(x, legend_top, Inches(1.2), Inches(0.25))
        tf = txBox.text_frame
        tf.paragraphs[0].text = f"■ {label}"
        tf.paragraphs[0].font.size = Pt(8)
        tf.paragraphs[0].font.color.rgb = color
