from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

COLORS = [
    RGBColor(0x44, 0x72, 0xC4),  # Blue
    RGBColor(0xED, 0x7D, 0x31),  # Orange
    RGBColor(0x70, 0xAD, 0x47),  # Green
    RGBColor(0xFF, 0xC0, 0x00),  # Gold
    RGBColor(0x5B, 0x9B, 0xD5),  # Light Blue
    RGBColor(0xA5, 0xA5, 0xA5),  # Gray
]


def create_presentation(title: str) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    return prs


def add_rect_node(slide, left, top, width, height, text, color, font_size=Pt(9)):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    shape.line.width = Pt(0.5)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    return shape


def add_connector(slide, start_left, start_top, end_left, end_top, label=""):
    connector = slide.shapes.add_connector(
        1,  # straight connector
        int(start_left), int(start_top),
        int(end_left), int(end_top),
    )
    connector.line.color.rgb = RGBColor(0x66, 0x66, 0x66)
    connector.line.width = Pt(1.5)
    return connector


def save_presentation(prs: Presentation, output_path: Path) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path
