from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from src.models.project import ProjectInfo
from src.generators.pptx_generator import (
    create_presentation, save_presentation, COLORS,
)
from src.i18n import t


def generate_aa07(project: ProjectInfo, output_dir: Path, locale: str = "zh") -> Path:
    """Generate AA-07 Application Architecture Diagram."""
    prs = create_presentation(t("pptx.arch_diagram", locale, name=project.name))
    slide = prs.slides[0]

    modules = project.modules
    if not modules:
        return save_presentation(prs, output_dir / t("file.aa07", locale))

    hierarchy = _build_hierarchy(project, locale)

    page_left = Inches(0.3)
    page_top = Inches(1.0)
    page_width = Inches(12.7)
    page_height = Inches(6.2)

    _draw_domain(slide, page_left, page_top, page_width, page_height, project.name, hierarchy, locale)

    return save_presentation(prs, output_dir / t("file.aa07", locale))


def _build_hierarchy(project: ProjectInfo, locale: str = "zh") -> dict:
    """Build: {app_group_name: {level1_module: [level2_classes]}}."""
    hierarchy = {}
    for module in project.modules:
        group_name = module.name
        level1_modules = {}

        if module.controllers:
            level1_modules[t("layer.api", locale)] = [c for c in module.controllers]
        if module.services:
            level1_modules[t("layer.business", locale)] = [s for s in module.services]
        if module.repositories:
            level1_modules[t("layer.data_access", locale)] = [r for r in module.repositories]
        if module.entities:
            level1_modules[t("layer.entity", locale)] = [e for e in module.entities]

        if not level1_modules:
            level1_modules[group_name] = []

        hierarchy[group_name] = level1_modules

    return hierarchy


def _draw_domain(slide, left, top, width, height, domain_name, hierarchy, locale: str = "zh"):
    """Draw the outermost domain container with all nested content."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.background()
    shape.line.color.rgb = COLORS[0]
    shape.line.width = Pt(2.5)

    title_height = Inches(0.45)
    title_shape = slide.shapes.add_shape(1, left, top, width, title_height)
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLORS[0]
    title_shape.line.color.rgb = COLORS[0]
    tf = title_shape.text_frame
    tf.paragraphs[0].text = t("pptx.domain_label", locale, name=domain_name)
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    content_left = left + Inches(0.15)
    content_top = top + title_height + Inches(0.15)
    content_width = width - Inches(0.3)
    content_height = height - title_height - Inches(0.3)

    groups = list(hierarchy.items())
    if not groups:
        return

    max_cols = min(len(groups), 4)
    num_rows = (len(groups) + max_cols - 1) // max_cols
    gap_x = Inches(0.12)
    gap_y = Inches(0.12)

    group_width = (content_width - gap_x * (max_cols - 1)) / max_cols
    group_height = (content_height - gap_y * (num_rows - 1)) / num_rows

    for idx, (group_name, level1_modules) in enumerate(groups):
        col = idx % max_cols
        row = idx // max_cols

        g_left = content_left + (group_width + gap_x) * col
        g_top = content_top + (group_height + gap_y) * row

        color = COLORS[(idx + 1) % len(COLORS)]
        _draw_app_group(slide, g_left, g_top, group_width, group_height, group_name, level1_modules, color)


def _draw_app_group(slide, left, top, width, height, group_name, level1_modules, color):
    """Draw an app group box with its level-1 modules inside."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.background()
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)

    title_height = Inches(0.35)
    title_shape = slide.shapes.add_shape(1, left, top, width, title_height)
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = color
    title_shape.line.color.rgb = color
    tf = title_shape.text_frame
    tf.paragraphs[0].text = group_name
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    mod_left = left + Inches(0.08)
    mod_top = top + title_height + Inches(0.08)
    mod_width = width - Inches(0.16)
    mod_height_total = height - title_height - Inches(0.16)

    level1_items = list(level1_modules.items())
    if not level1_items:
        return

    gap = Inches(0.06)
    single_height = (mod_height_total - gap * (len(level1_items) - 1)) / len(level1_items)
    single_height = max(single_height, Inches(0.3))

    for l1_idx, (l1_name, l2_classes) in enumerate(level1_items):
        l1_top = mod_top + (single_height + gap) * l1_idx
        _draw_level1_module(slide, mod_left, l1_top, mod_width, single_height, l1_name, l2_classes, color)


def _draw_level1_module(slide, left, top, width, height, l1_name, l2_classes, parent_color):
    """Draw a level-1 module row with level-2 items inside."""
    light_color = _lighten_color(parent_color)

    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = light_color
    shape.line.color.rgb = parent_color
    shape.line.width = Pt(0.75)

    label_width = Inches(0.7)
    label_shape = slide.shapes.add_textbox(left + Inches(0.04), top + Inches(0.02), label_width, height - Inches(0.04))
    tf = label_shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = l1_name
    tf.paragraphs[0].font.size = Pt(7)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    if not l2_classes:
        return

    tag_left = left + label_width + Inches(0.08)
    tag_top = top + Inches(0.04)
    available_width = width - label_width - Inches(0.2)

    tag_height = Inches(0.22)
    tag_gap_x = Inches(0.06)
    tag_gap_y = Inches(0.04)
    current_x = tag_left
    current_y = tag_top

    for cls_name in l2_classes:
        tag_width = Inches(max(0.6, min(1.8, len(cls_name) * 0.07)))

        if current_x + tag_width > left + width - Inches(0.08):
            current_x = tag_left
            current_y += tag_height + tag_gap_y

        if current_y + tag_height > top + height - Inches(0.04):
            overflow_box = slide.shapes.add_textbox(current_x, current_y, Inches(0.3), tag_height)
            tf = overflow_box.text_frame
            tf.paragraphs[0].text = "..."
            tf.paragraphs[0].font.size = Pt(6)
            break

        tag_shape = slide.shapes.add_shape(1, current_x, current_y, tag_width, tag_height)
        tag_shape.fill.solid()
        tag_shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tag_shape.line.color.rgb = parent_color
        tag_shape.line.width = Pt(0.5)

        tf = tag_shape.text_frame
        tf.word_wrap = False
        tf.paragraphs[0].text = cls_name
        tf.paragraphs[0].font.size = Pt(6)
        tf.paragraphs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        current_x += tag_width + tag_gap_x


def _lighten_color(color: RGBColor) -> RGBColor:
    """Create a lighter version of a color for backgrounds."""
    r = min(255, color[0] + (255 - color[0]) * 80 // 100)
    g = min(255, color[1] + (255 - color[1]) * 80 // 100)
    b = min(255, color[2] + (255 - color[2]) * 80 // 100)
    return RGBColor(r, g, b)
