import pytest
from pathlib import Path
from openpyxl import load_workbook
from pptx import Presentation

from src.scanner.project_scanner import scan_project
from src.analyzers.module_analyzer import analyze_modules
from src.analyzers.api_analyzer import analyze_apis
from src.analyzers.integration_analyzer import analyze_integrations
from src.generators.app_architecture_generator import generate_app_architecture
from src.generators.aa07_generator import generate_aa07
from src.generators.aa08_generator import generate_aa08


@pytest.fixture
def sample_project():
    proj_path = Path(__file__).parent / "fixtures" / "sample-java-project"
    project = scan_project(proj_path)
    project = analyze_modules(project)
    return project


@pytest.fixture
def sample_endpoints(sample_project):
    return analyze_apis(sample_project)


@pytest.fixture
def sample_integrations(sample_project):
    return analyze_integrations(sample_project)


def test_english_excel_headers(sample_project, sample_endpoints, sample_integrations, tmp_path):
    excel_path = generate_app_architecture(
        sample_project, sample_endpoints, sample_integrations, tmp_path, locale="en"
    )
    wb = load_workbook(excel_path)

    # AA-01 headers
    ws = wb.worksheets[0]
    headers = [cell.value for cell in ws[1]]
    assert headers[0] == "Domain ID"
    assert headers[-1] == "Build Status"
    assert "应用" not in str(headers)

    # AA-02 headers
    ws = wb.worksheets[1]
    headers = [cell.value for cell in ws[1]]
    assert "Function ID" in headers
    assert "HTTP Method" in headers

    # AA-05 headers
    ws = wb.worksheets[4]
    headers = [cell.value for cell in ws[1]]
    assert "Source Module" in headers
    assert "Target Module" in headers


def test_english_excel_cell_values(sample_project, sample_endpoints, sample_integrations, tmp_path):
    excel_path = generate_app_architecture(
        sample_project, sample_endpoints, sample_integrations, tmp_path, locale="en"
    )
    wb = load_workbook(excel_path)

    # AA-01: "Built" instead of "已建", layer names in English
    ws = wb.worksheets[0]
    statuses = [row[-1] for row in ws.iter_rows(min_row=2, values_only=True)]
    assert all(s == "Built" for s in statuses)

    layer_names = [row[5] for row in ws.iter_rows(min_row=2, values_only=True)]
    assert "API Service" in layer_names or "Business Logic" in layer_names
    assert "接口服务" not in layer_names

    # AA-04: "Microservice Application" instead of "微服务应用"
    ws = wb.worksheets[3]
    categories = [row[3] for row in ws.iter_rows(min_row=2, values_only=True)]
    assert all(c == "Microservice Application" for c in categories)

    # AA-05: "Yes"/"No" instead of "是"/"否"
    ws = wb.worksheets[4]
    cross_domain = [row[-1] for row in ws.iter_rows(min_row=2, values_only=True)]
    assert all(v in ("Yes", "No") for v in cross_domain)


def test_english_excel_sheet_names(sample_project, sample_endpoints, sample_integrations, tmp_path):
    excel_path = generate_app_architecture(
        sample_project, sample_endpoints, sample_integrations, tmp_path, locale="en"
    )
    wb = load_workbook(excel_path)
    for name in wb.sheetnames:
        assert "AA-0" in name
        assert "应用" not in name


def test_english_aa07_pptx(sample_project, tmp_path):
    aa07_path = generate_aa07(sample_project, tmp_path, locale="en")
    assert aa07_path.name == "AA-07_Application_Architecture.pptx"

    prs = Presentation(str(aa07_path))
    slide = prs.slides[0]

    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                texts.append(text)

    assert any("Application Architecture" in t for t in texts)
    assert any("Domain:" in t for t in texts)
    assert "应用架构图" not in str(texts)
    assert "应用域" not in str(texts)


def test_english_aa08_pptx(sample_project, sample_integrations, tmp_path):
    aa08_path = generate_aa08(sample_project, sample_integrations, tmp_path, locale="en")
    assert aa08_path.name == "AA-08_Application_Integration.pptx"

    prs = Presentation(str(aa08_path))
    slide = prs.slides[0]

    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                texts.append(text)

    assert any("Application Integration" in t for t in texts)
    assert "应用集成图" not in str(texts)


def test_chinese_output_unchanged(sample_project, sample_endpoints, sample_integrations, tmp_path):
    """Verify default locale='zh' produces same output as before."""
    excel_path = generate_app_architecture(
        sample_project, sample_endpoints, sample_integrations, tmp_path
    )
    wb = load_workbook(excel_path)

    ws = wb.worksheets[0]
    headers = [cell.value for cell in ws[1]]
    assert headers[0] == "应用域编号"

    statuses = [row[-1] for row in ws.iter_rows(min_row=2, values_only=True)]
    assert all(s == "已建" for s in statuses)
